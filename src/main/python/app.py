import os
import uuid
import json
import time
from flask import Flask, request, jsonify, session, redirect, url_for, render_template, make_response, send_file
from dotenv import load_dotenv
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
import google.generativeai as genai
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime, timezone, timedelta
from authlib.integrations.flask_client import OAuth
from werkzeug.utils import secure_filename
from utils import get_image_description, generate_image_from_prompt
import re
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image

# --- PATH & ENV CONFIGURATION ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(BASE_DIR, '../../../'))
load_dotenv(os.path.join(PROJECT_ROOT, '.env'))
os.environ['OAUTHLIB_INSECURE_TRANSPORT'] = '1'

template_dir = os.path.join(PROJECT_ROOT, 'templates')
static_dir = os.path.join(PROJECT_ROOT, 'static')
UPLOAD_FOLDER = os.path.join(static_dir, 'uploads')
PRESENTATION_OUTPUT_FOLDER = os.path.join(static_dir, 'presentations')

# --- GLOBAL CONFIGURATIONS ---
# Models list: Prioritizes gemini-2.5-flash, falls back to 1.5 versions if needed.
GEMINI_MODELS_TO_TRY = [
    'gemini-2.5-flash',     # Primary: Best price-performance & speed
    'gemini-1.5-flash',     # Fallback 1: Very stable high rate limits
    'gemini-1.5-flash-001', # Fallback 2: Specific version tag
    'gemini-1.5-pro'        # Fallback 3: Advanced reasoning
]

DEEP_RESEARCH_LIMIT = 5
DEEP_RESEARCH_MAX_RETRIES = 1
DEEP_RESEARCH_RETRY_DELAY = 2
DEEP_SEARCH_RESULTS_COUNT = 8
CANVAS_SUB_QUERY_COUNT = 4

genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

app = Flask(__name__, template_folder=template_dir, static_folder=static_dir)

app.config['SECRET_KEY'] = os.getenv('FLASK_SECRET_KEY')
app.config['SQLALCHEMY_DATABASE_URI'] = os.getenv('DATABASE_URL')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
db = SQLAlchemy(app)
oauth = OAuth(app)

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(PRESENTATION_OUTPUT_FOLDER, exist_ok=True)

google = oauth.register(
    name='google',
    client_id=os.getenv("GOOGLE_CLIENT_ID"),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
    access_token_url='https://accounts.google.com/o/oauth2/token',
    access_token_params=None,
    authorize_url='https://accounts.google.com/o/oauth2/auth',
    authorize_params=None,
    api_base_url='https://www.googleapis.com/oauth2/v1/',
    userinfo_endpoint='https://openidconnect.googleapis.com/v1/userinfo',
    client_kwargs={'scope': 'openid email profile'},
    jwks_uri='https://www.googleapis.com/oauth2/v3/certs'
)

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    google_id = db.Column(db.String(30), unique=True, nullable=False)
    name = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(100), unique=True, nullable=False)
    profile_pic = db.Column(db.String(200), nullable=True)
    conversations = db.relationship('Conversation', backref='user', lazy=True, cascade="all, delete-orphan")
    image_generations_count = db.Column(db.Integer, nullable=False, default=0)
    last_generation_timestamp = db.Column(db.DateTime, nullable=True)
    deep_research_count = db.Column(db.Integer, nullable=False, default=0)
    last_deep_research_timestamp = db.Column(db.DateTime, nullable=True)
    presentation_slides = db.relationship('PresentationSlide', backref='user', lazy=True, cascade="all, delete-orphan")

class Conversation(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(150), nullable=False)
    timestamp = db.Column(db.DateTime, nullable=False, default=lambda: datetime.now(timezone.utc))
    messages = db.relationship('ChatMessage', backref='conversation', lazy=True, cascade="all, delete-orphan")
    presentation_slides = db.relationship('PresentationSlide', backref='conversation', lazy=True, cascade="all, delete-orphan")

class ChatMessage(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    conversation_id = db.Column(db.Integer, db.ForeignKey('conversation.id'), nullable=False)
    question = db.Column(db.Text, nullable=False)
    answer = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, nullable=False, default=lambda: datetime.now(timezone.utc))
    sources = db.Column(db.JSON, nullable=True)
    canvas_data = db.Column(db.JSON, nullable=True)

class PresentationSlide(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    conversation_id = db.Column(db.Integer, db.ForeignKey('conversation.id'), nullable=False)
    content = db.Column(db.Text, nullable=False)
    slide_order = db.Column(db.Integer, nullable=False)
    timestamp = db.Column(db.DateTime, nullable=False, default=lambda: datetime.now(timezone.utc))
    chat_message_id = db.Column(db.Integer, db.ForeignKey('chat_message.id'), nullable=True)

# --- HELPER: ROBUST GEMINI CALLER ---
def generate_content_with_fallback(prompt, history=None):
    """
    Tries multiple Gemini models starting with gemini-2.5-flash.
    Handles 429 Quota errors and 404 Not Found errors gracefully.
    """
    last_error = None
    
    for model_name in GEMINI_MODELS_TO_TRY:
        try:
            model = genai.GenerativeModel(model_name)
            if history:
                chat_session = model.start_chat(history=history)
                response = chat_session.send_message(prompt)
            else:
                response = model.generate_content(prompt)
            
            # If successful, return the text
            return response.text
            
        except Exception as e:
            error_str = str(e).lower()
            print(f"[DEBUG] Model {model_name} failed: {error_str}")
            
            # If it's a quota error (429) or Not Found (404), wait briefly then try next model
            if "429" in error_str or "quota" in error_str or "404" in error_str or "not found" in error_str:
                time.sleep(1)
            
            last_error = e
            continue # Try next model in the list
            
    # If all models fail, raise the last error
    raise last_error

def classify_query(query: str):
    try:
        prompt = f"""Analyze the user's query and classify its intent. Respond with only a single word: 'SEARCH', 'CONVERSATIONAL', or 'GENERATIVE'.
        - 'SEARCH': For questions requiring factual, real-time web info (e.g., "What's the weather today?", "Latest news on X", "Who won the World Cup 2022?").
        - 'CONVERSATIONAL': For greetings, small talk, or direct commands that don't require external data (e.g., "Hello", "How are you?", "Tell me a joke").
        - 'GENERATIVE': For knowledge-based questions that the AI can answer from its training data, without needing a web search (e.g., "Explain quantum computing", "What is the capital of France?", "Summarize the history of AI").
        User Query: "{query}"
        Classification:"""
        
        response_text = generate_content_with_fallback(prompt)
        classification = response_text.strip().upper()
        
        if "SEARCH" in classification: return "SEARCH"
        if "GENERATIVE" in classification: return "GENERATIVE"
        if re.match(r"^\s*(hello|hi|hey|greetings|hru|how are you|yo)\s*$", query.lower()):
            return "CONVERSATIONAL_GREETING"
        return "CONVERSATIONAL"
    except Exception as e:
        print(f"Error classifying query (Fallback): {e}")
        return "GENERATIVE" 

def search_the_web(query: str):
    try:
        service = build("customsearch", "v1", developerKey=os.getenv("GOOGLE_API_KEY"))
        result = service.cse().list(q=query, cx=os.getenv("GOOGLE_CSE_ID"), num=5).execute()
        return result.get('items', [])
    except Exception as e:
        print(f"Web search error: {e}")
        return []

def search_youtube(query: str):
    try:
        service = build("youtube", "v3", developerKey=os.getenv("GOOGLE_API_KEY"))
        request = service.search().list(q=query, part='snippet', maxResults=3, type='video')
        response = request.execute()
        return response.get('items', [])
    except Exception as e:
        print(f"YouTube search error: {e}")
        return []

def search_images(query: str, num_images=2):
    try:
        service = build("customsearch", "v1", developerKey=os.getenv("GOOGLE_API_KEY"))
        result = service.cse().list(q=query, cx=os.getenv("GOOGLE_CSE_ID"), num=num_images, searchType='image', safe='high').execute()
        return result.get('items', [])
    except Exception as e:
        print(f"Image search error: {e}")
        return []

def deep_research_the_web(query: str):
    google_sources = []
    try:
        cse_id = os.getenv("GOOGLE_DEEP_SEARCH_CSE_ID")
        if not cse_id:
            print("Error: GOOGLE_DEEP_SEARCH_CSE_ID not found in .env file.")
        else:
            service = build("customsearch", "v1", developerKey=os.getenv("GOOGLE_API_KEY"))
            result = service.cse().list(q=query, cx=cse_id, num=DEEP_SEARCH_RESULTS_COUNT).execute()

            for item in result.get('items', []):
                google_sources.append({
                    "url": item.get('link', ''),
                    "title": item.get('title', ''),
                    "snippet": item.get('snippet', 'No snippet available.'),
                    "source_type": "Google Academic Search"
                })
    except Exception as e:
        print(f"Google Deep Search error: {e}")

    return google_sources

def get_chat_history_for_api(conversation_id):
    history_for_api = []
    if conversation_id:
        previous_messages = ChatMessage.query.filter_by(conversation_id=conversation_id).order_by(ChatMessage.timestamp.asc()).all()
        for msg in previous_messages:
            if not msg.question.startswith(("[IMAGE:", "[IMAGE_SEARCH]", "[DEEP_RESEARCH]")):
                history_for_api.append({'role': 'user', 'parts': [{'text': msg.question}]})
                history_for_api.append({'role': 'model', 'parts': [{'text': msg.answer}]})
            elif msg.question.startswith("[REPORT CONTEXT]"):
                history_for_api.append({'role': 'model', 'parts': [{'text': msg.answer}]})

    return history_for_api

def generate_answer(context: str, question: str, history: list = None, use_internal_knowledge: bool = False):
    """
    Generates an answer using Gemini with dynamic tone adjustment.
    - If Context is present (Search Results) -> Professional & Structured.
    - If Internal Knowledge (Common Sense) -> Friendly & Conversational.
    """
    try:
        language_instruction = f"CRITICAL INSTRUCTION: Your ENTIRE response MUST be in the same language as the user's question. The user's question is: '{question}'."

        # --- DYNAMIC TONE ADJUSTMENT LOGIC ---
        if context:
            # SEARCH MODE: Be professional, objective, and cite sources.
            tone_instruction = """
            **TONE & STYLE:** - Adopt a **Professional, Objective, and Informative** tone.
            - Structure the answer clearly with headings or bullet points where appropriate.
            - Strictly cite the provided sources.
            """
            prompt = f"""You are Synapse AI Chat, a high-precision research assistant. 
            {tone_instruction}
            
            **RULES:**
            1. **NEVER USE YOUR INTERNAL KNOWLEDGE** for the core facts; strictly use the provided context.
            2. **CITE SOURCES:** End sentences with citations like [1], [2].
            3. **LANGUAGE:** {language_instruction}
            
            **USER QUESTION:** "{question}"
            
            **WEB SEARCH CONTEXT:**
            ---
            {context}
            ---
            Based **only** on the text above, provide your answer."""
            
        elif use_internal_knowledge:
            # GENERATIVE / CONVERSATIONAL MODE: Be friendly, helpful, and natural.
            tone_instruction = """
            **TONE & STYLE:** - Adopt a **Friendly, Conversational, and Helpful** tone.
            - Be direct and easy to understand.
            - If it's a greeting, reply warmly. If it's a complex question, explain it simply.
            """
            prompt = f"""You are Synapse AI Chat, a helpful and professional AI assistant created by 'H7'.
            {tone_instruction}
            
            Engage in common-sense conversation. Respond comprehensively using your knowledge.
            {language_instruction}
            
            User's current message: '{question}'."""
            
        else:
            # Fallback
            prompt = f"You are Synapse AI Chat. Answer the user's question: '{question}'. {language_instruction}"
        
        return generate_content_with_fallback(prompt, history=history)

    except Exception as e:
        print(f"Error generating answer: {e}")
        return "I seem to be having some technical difficulties with the AI service. Please try again later."

def generate_deep_research_report(query: str):
    try:
        academic_sources = deep_research_the_web(query)
        image_sources = search_images(query, num_images=2)

        context = ""
        sources_for_report = []
        if academic_sources:
            context_parts = []
            for item in academic_sources:
                source_url = item.get('link', '')
                sources_for_report.append({"url": source_url, "title": item.get('title', ''), "snippet": item.get('snippet', '')})
                context_parts.append(f"Source [{len(sources_for_report)}]: Type: Google Academic Source\nTitle: {item.get('title', '')}\nContent: {item.get('snippet', '')}\nURL: {source_url}")
            context = "\n---\n".join(context_parts)

        images_for_report = [{"src": img.get("link"), "alt": img.get("title")} for img in image_sources if img.get("link")]

        prompt_academic_context_message = context if context else "No academic sources were found for this query. Generate a report acknowledging this limitation and provide general information if possible, otherwise state that specific academic references are unavailable."
        sources_for_report_json = json.dumps(sources_for_report) if sources_for_report else "[]"

        prompt = f"""
        You are a professional AI research synthesizer. Generate a comprehensive, structured academic report as a single, valid JSON object.
        **TOPIC:** "{query}"
        **CONTEXT:**
        ---
        {prompt_academic_context_message}
        ---
        **CRITICAL INSTRUCTIONS:**
        1.  **JSON Structure:** Your entire response MUST be a single, valid JSON object following this exact professional structure. Ensure all string values like "content", "abstract" etc. are properly formatted JSON strings. If no information is available for a field, provide a short explanatory string, NOT null.
            {{
                "report_title": "A comprehensive, academic title for the research report",
                "abstract": "A concise, one-paragraph summary of the entire report, including key findings and conclusions.",
                "sections": [
                    {{"heading": "Introduction", "content": "A detailed introduction providing background on the topic. Cite sources heavily using [1], [2], etc."}},
                    {{"heading": "Key Findings and Analysis", "content": "Synthesize the main arguments from the sources into several detailed paragraphs. Insert image placeholders like '[IMAGE_1]' where most relevant."}},
                    {{"heading": "Discussion and Implications", "content": "Discuss the significance and implications of the findings. Analyze the information already presented. Cite sources."}}
                ],
                "images": {json.dumps(images_for_report)},
                "key_concepts": ["A list of 3-5 key technical terms or concepts."],
                "source_summary": [{{"source_number": 1, "summary": "A one-sentence summary of what Source [1] contributed."}}],
                "conclusion": "A strong concluding paragraph summarizing the main points and suggesting areas for future research.",
                "references": {sources_for_report_json}
            }}
        2.  **Professional Tone & Synthesis:** Write in a formal, academic tone. Synthesize information to create a cohesive narrative. Avoid redundancy.
        3.  **Image Integration:** In the section content, place `[IMAGE_1]`, `[IMAGE_2]`.
        4.  **Language:** Write in the same language as the user's topic: "{query}".

        Generate the complete JSON object now.
        """
        
        response_text = generate_content_with_fallback(prompt)
        if not response_text:
            raise ValueError("Received empty response from generative model.")

        print(f"\n--- [DEBUG] Raw AI response for Deep Research: ---\n{response_text}\n--- END RAW AI RESPONSE ---")

        cleaned_response = response_text.strip().replace("```json", "").replace("```", "")

        report_data = json.loads(cleaned_response)

        final_report_data = {
            'report_title': report_data.get('report_title', 'Deep Research Report (Untitled)'),
            'abstract': report_data.get('abstract', 'No abstract provided.'),
            'sections': report_data.get('sections', []),
            'images': report_data.get('images', []),
            'key_concepts': report_data.get('key_concepts', []),
            'source_summary': report_data.get('source_summary', []),
            'conclusion': report_data.get('conclusion', 'No conclusion provided.'),
            'references': report_data.get('references', [])
        }

        if 'sections' in final_report_data and isinstance(final_report_data['sections'], list):
            for i, section in enumerate(final_report_data['sections']):
                if not isinstance(section, dict):
                    final_report_data['sections'][i] = {"heading": f"Section {i+1}", "content": "Content unavailable."}
                else:
                    section['heading'] = section.get('heading', f"Section {i+1}")
                    section['content'] = section.get('content', "Content unavailable.")

        return final_report_data

    except json.JSONDecodeError as e:
        print(f"JSON Parsing Error in generate_deep_research_report: {e}")
        return {"error": f"Failed to parse AI response as JSON. Raw response might be malformed.", "raw_ai_response": "Error"}
    except Exception as e:
        print(f"Error generating or parsing deep research report: {e}")
        return {"error": f"An unexpected error occurred during report generation: {e}"}

def download_image_for_pptx(image_url):
    try:
        response = requests.get(image_url, stream=True, timeout=10)
        response.raise_for_status()
        img_data = BytesIO(response.content)

        try:
            img = Image.open(img_data)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            output_io = BytesIO()
            img.save(output_io, format='PNG')
            output_io.seek(0)
            return output_io
        except Exception as img_e:
            print(f"PIL image processing failed for {image_url}: {img_e}. Attempting to return raw data.")
            img_data.seek(0)
            return img_data
    except requests.exceptions.RequestException as e:
        print(f"Error downloading image {image_url}: {e}")
        return None
    except Exception as e:
        print(f"Unexpected error in download_image_for_pptx: {e}")
        return None

def get_synthesis_sub_queries(main_topic: str):
    try:
        prompt = f"""You are an AI assistant specialized in breaking down complex topics into core sub-questions.
        Given the main topic "{main_topic}", generate exactly {CANVAS_SUB_QUERY_COUNT} distinct, concise, and highly relevant sub-questions that are crucial for a comprehensive understanding of the main topic.
        Respond with a JSON array of strings, where each string is a sub-question.
        Example: ["Sub-question 1", "Sub-question 2", ...]
        """
        response_text = generate_content_with_fallback(prompt)
        cleaned_response = response_text.strip().replace("```json", "").replace("```", "")
        sub_queries = json.loads(cleaned_response)

        if not isinstance(sub_queries, list) or not all(isinstance(q, str) for q in sub_queries):
            raise ValueError("Sub-query generation did not return a list of strings.")

        return sub_queries
    except Exception as e:
        print(f"Error generating synthesis sub-queries for '{main_topic}': {e}")
        return [main_topic] # Fallback

def run_multiple_deep_searches(sub_queries: list):
    all_reports = []
    for sub_query in sub_queries:
        print(f"[DEBUG] Running deep research for sub-query: {sub_query}")
        report_json = {"error": "Failed to generate report for sub-query."}
        try:
            report_json = generate_deep_research_report(sub_query)
            if "error" not in report_json:
                all_reports.append(report_json)
            else:
                print(f"[DEBUG] Failed for sub-query '{sub_query}': {report_json.get('error')}.")
        except Exception as e:
            print(f"[DEBUG] Exception for sub-query '{sub_query}': {e}.")
    return all_reports

def transform_report_to_canvas_data(report_json):
    nodes = []
    edges = []
    node_id_counter = 1

    if not report_json or "error" in report_json:
        return {"nodes": [], "edges": []}

    main_topic_id = node_id_counter
    nodes.append({
        "id": main_topic_id,
        "label": report_json.get('report_title', 'Research Topic'),
        "color": '#3b82f6',
        "shape": 'dot',
        "size": 30,
        "font": {"size": 20, "color": '#e5e7eb'}
    })
    node_id_counter += 1

    key_concept_ids = {}
    key_concepts = report_json.get('key_concepts', [])
    concept_colors = ['#22c55e', '#eab308', '#f43f5e', '#8b5cf6', '#14b8a6']
    for i, concept in enumerate(key_concepts):
        concept_id = node_id_counter
        key_concept_ids[concept.lower()] = concept_id
        nodes.append({
            "id": concept_id,
            "label": concept,
            "color": concept_colors[i % len(concept_colors)],
            "size": 20
        })
        edges.append({"from": main_topic_id, "to": concept_id})
        node_id_counter += 1

    sections = report_json.get('sections', [])
    for section in sections:
        section_id = node_id_counter
        nodes.append({
            "id": section_id,
            "label": section.get('heading', 'Section'),
            "shape": 'ellipse',
            "color": '#1f2937'
        })
        edges.append({"from": main_topic_id, "to": section_id})
        node_id_counter += 1

        finding_id = node_id_counter
        finding_content = section.get('content', 'No content.')
        tooltip_text = (finding_content[:300] + '...') if len(finding_content) > 300 else finding_content
        nodes.append({
            "id": finding_id,
            "label": 'Finding',
            "shape": 'box',
            "title": tooltip_text
        })
        edges.append({"from": section_id, "to": finding_id})
        node_id_counter += 1

    source_summaries = report_json.get('source_summary', [])
    if source_summaries:
        summary_hub_id = node_id_counter
        nodes.append({"id": summary_hub_id, "label": "Sources", "color": "#6b7280", "size": 18})
        edges.append({"from": main_topic_id, "to": summary_hub_id})
        node_id_counter += 1

        for summary in source_summaries:
            source_id = node_id_counter
            source_label = f"Source [{summary.get('source_number', '?')}]"
            source_title = summary.get('summary', 'No summary available.')
            nodes.append({
                "id": source_id,
                "label": source_label,
                "shape": 'database',
                "title": source_title
            })
            edges.append({"from": summary_hub_id, "to": source_id})
            node_id_counter += 1

    return {"nodes": nodes, "edges": []}

def add_hyperlink(paragraph, url, text, style='Hyperlink'):
    if paragraph is None or not hasattr(paragraph, 'part'):
        print(f"[ERROR] Invalid paragraph object passed to add_hyperlink: {paragraph}")
        if paragraph and hasattr(paragraph, 'add_run'):
            paragraph.add_run(text)
        return None

    part = paragraph.part
    if not url:
        paragraph.add_run(text)
        return None

    r_id = part.relate_to(url, RT.HYPERLINK)
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)

    new_run = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')

    if style and style in paragraph.part.document.styles:
        rStyle = OxmlElement('w:rStyle')
        rStyle.set(qn('w:val'), style)
        rPr.append(rStyle)
    else:
        color = OxmlElement('w:color')
        color.set(qn('w:val'), '0000FF')
        rPr.append(color)
        underline = OxmlElement('w:u')
        underline.set(qn('w:val'), 'single')
        rPr.append(underline)

    new_run.append(rPr)
    t = OxmlElement('w:t')
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._element.append(hyperlink)
    return hyperlink

def add_markdown_content_to_docx(document, text_content, images_for_report, references):
    if text_content is None: return

    for line in text_content.split('\n'):
        line = line.strip()
        if not line: continue

        if line.startswith("## "):
            document.add_heading(line[3:], level=2)
            continue
        if line.startswith("* ") or line.startswith("- "):
            document.add_paragraph(line[2:], style='List Bullet')
            continue

        document.add_paragraph(line)

@app.route('/export-to-word', methods=['POST'])
def export_to_word_route():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401

    try:
        report_data = request.get_json()

        if not report_data:
            return jsonify({"error": "No report data provided."}), 400

        if "error" in report_data:
            return jsonify({"error": f"Cannot export a failed report: {report_data['error']}"}), 400

        document = Document()

        document.add_heading(report_data.get('report_title', 'Deep Research Report'), level=0).alignment = WD_ALIGN_PARAGRAPH.CENTER
        document.add_paragraph(f"\nReport Generated by Synapse AI Chat", style='Subtitle').alignment = WD_ALIGN_PARAGRAPH.CENTER
        document.add_paragraph(f"Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}", style='Normal').alignment = WD_ALIGN_PARAGRAPH.CENTER
        document.add_page_break()

        images = report_data.get('images', [])
        references = report_data.get('references', [])

        document.add_heading('Abstract', level=1)
        add_markdown_content_to_docx(document, report_data.get('abstract', ''), images, references)

        if report_data.get('sections'):
            for i, section in enumerate(report_data['sections']):
                document.add_heading(section.get('heading', 'Section'), level=1)
                add_markdown_content_to_docx(document, section.get('content', ''), images, references)

        if report_data.get('key_concepts'):
            document.add_heading('Key Concepts', level=1)
            for concept in report_data.get('key_concepts', []):
                document.add_paragraph(concept, style='List Bullet')

        if report_data.get('source_summary'):
            document.add_heading('Source Summaries', level=1)
            for summary in report_data.get('source_summary', []):
                document.add_paragraph(f"Source [{summary.get('source_number')}]: {summary.get('summary')}", style='Normal')

        document.add_heading('Conclusion', level=1)
        add_markdown_content_to_docx(document, report_data.get('conclusion', ''), images, references)

        document.add_page_break()
        document.add_heading('References', level=1)
        if references:
            for i, ref in enumerate(references):
                p = document.add_paragraph(style='List Number')
                p.add_run(f"[{i+1}] ").bold = True
                p.add_run(ref.get('title', 'Untitled'))
                snippet = ref.get('snippet', 'No snippet provided.')
                if i > 0:
                    sep_para = document.add_paragraph()
                    sep_run = sep_para.add_run('—'*50)
                    sep_run.font.size = Pt(5)
                p_snippet = document.add_paragraph(f"Source: {snippet}", style='Normal')
                p_snippet.paragraph_format.left_indent = Inches(0.2)
        else:
            document.add_paragraph("No references available.")

        file_stream = BytesIO()
        document.save(file_stream)
        file_stream.seek(0)
        report_title_safe = re.sub(r'[\\/:*?"<>|]', '', report_data.get('report_title', 'Report'))
        file_name = f"{report_title_safe.replace(' ', '_')}.docx"

        response = make_response(file_stream.read())
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        response.headers['Content-Disposition'] = f'attachment; filename="{file_name}"'
        return response

    except Exception as e:
        print(f"CRITICAL Error in /export-to-word route: {e}")
        return jsonify({"error": f"Failed to export Word document due to an unexpected server error: {str(e)}"}), 500

@app.route('/')
def home():
    if 'user' in session: return redirect(url_for('chat_page'))
    return render_template('index.html')

@app.route('/chat_page')
@app.route('/chat_page/new')
@app.route('/chat_page/<int:conv_id>')
def chat_page(conv_id=None):
    if 'user' not in session: return redirect(url_for('home'))
    return render_template('chat.html', user=session.get('user'), initial_conv_id=conv_id)

@app.route('/settings')
def settings_page():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    return render_template('settings.html', user=session.get('user'))

@app.route('/login')
def login():
    redirect_uri = url_for('callback', _external=True)
    return google.authorize_redirect(redirect_uri)

@app.route('/callback')
def callback():
    token = google.authorize_access_token()
    user_info = token.get('userinfo')
    if not user_info or not user_info.get('sub'): return "Error: Could not fetch user information from Google.", 400
    google_id = user_info.get('sub')
    email = user_info.get('email')
    name = user_info.get('name') or email
    picture = user_info.get('picture')
    user = User.query.filter_by(google_id=google_id).first()
    if not user:
        user = User(google_id=google_id, name=name, email=email, profile_pic=picture)
        db.session.add(user)
    else:
        user.name = name
        user.email = email
        user.profile_pic = picture
    db.session.commit()
    session['user'] = {'id': user.id, 'name': user.name, 'picture': user.profile_pic, 'email': user.email}
    return redirect(url_for('chat_page'))

@app.route('/logout')
def logout():
    session.pop('user', None)
    return redirect('/')

@app.route('/history', methods=['GET'])
def history():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    conversations = Conversation.query.filter_by(user_id=user_id).order_by(Conversation.timestamp.desc()).all()
    return jsonify([{"id": conv.id, "title": conv.title} for conv in conversations])

@app.route('/history/delete_all', methods=['DELETE'])
def delete_all_history():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    try:
        Conversation.query.filter_by(user_id=user_id).delete()
        db.session.commit()
        return jsonify({"success": True, "message": "All chat history has been deleted."})
    except Exception as e:
        db.session.rollback(); print(f"Error deleting all history: {e}")
        return jsonify({"error": "An internal server error occurred."}), 500

@app.route('/account/delete', methods=['DELETE'])
def delete_account():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    try:
        user = db.session.get(User, user_id)
        if user:
            db.session.delete(user)
            db.session.commit()
            session.pop('user', None)
            return jsonify({"success": True, "redirect_url": url_for('home')})
        return jsonify({"error": "User not found."}), 404
    except Exception as e:
        db.session.rollback(); print(f"Error deleting account: {e}")
        return jsonify({"error": "An internal server error occurred."}), 500

@app.route('/conversation/<int:conv_id>', methods=['GET', 'PUT', 'DELETE'])
def manage_conversation(conv_id):
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    conversation = Conversation.query.filter_by(id=conv_id, user_id=session['user']['id']).first()
    if not conversation:
        return jsonify({"error": "Conversation not found or access denied"}), 404

    if request.method == 'GET':
        messages = ChatMessage.query.filter_by(conversation_id=conv_id).order_by(ChatMessage.timestamp.asc()).all()
        messages_data = []
        for msg in messages:
            msg_data = {"question": msg.question, "answer": msg.answer, "sources": msg.sources}
            if msg.canvas_data:
                msg_data["canvas_data"] = msg.canvas_data
            messages_data.append(msg_data)
        return jsonify(messages_data)
    elif request.method == 'PUT':
        data = request.get_json()
        if not data or 'title' not in data or not data['title'].strip(): return jsonify({"error": "New title is required"}), 400
        conversation.title = data['title'].strip()[:150]
        conversation.timestamp = datetime.now(timezone.utc)
        db.session.commit()
        return jsonify({"success": True, "message": "Conversation renamed."})
    elif request.method == 'DELETE':
        db.session.delete(conversation)
        db.session.commit()
        return jsonify({"success": True, "message": "Conversation deleted."})

def get_or_create_conversation(user_id, conv_id, title):
    if conv_id:
        conv = db.session.get(Conversation, conv_id)
        if conv and conv.user_id == user_id:
            return conv
    new_conv = Conversation(user_id=user_id, title=title)
    db.session.add(new_conv)
    db.session.commit()
    return new_conv

@app.route('/deep-research-status', methods=['GET'])
def deep_research_status():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    user = db.session.get(User, user_id)
    if not user: return jsonify({"error": "User not found."}), 404
    now = datetime.now(timezone.utc)
    limit = DEEP_RESEARCH_LIMIT
    limit_reached = False
    reset_time_iso = None

    if user.last_deep_research_timestamp and user.last_deep_research_timestamp.tzinfo is None:
        user.last_deep_research_timestamp = user.last_deep_research_timestamp.replace(tzinfo=timezone.utc)

    if user.last_deep_research_timestamp and (now - user.last_deep_research_timestamp) > timedelta(hours=24):
        user.deep_research_count = 0
        user.last_deep_research_timestamp = None
        db.session.commit()

    generations_left = limit - user.deep_research_count

    if generations_left <= 0:
        limit_reached = True
        if user.last_deep_research_timestamp:
            reset_time = user.last_deep_research_timestamp + timedelta(hours=24)
            reset_time_iso = reset_time.isoformat()

    return jsonify({
        "limit_reached": limit_reached,
        "generations_left": generations_left,
        "reset_time": reset_time_iso
    })

@app.route('/deep-research', methods=['POST'])
def deep_research_route():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    query = data.get('query')
    conversation_id = data.get('conversation_id')
    if not query: return jsonify({"error": "A research query is required."}), 400

    query_type = classify_query(query)
    # UPDATED: Allow deep research to proceed for GENERATIVE and SEARCH
    if query_type == "CONVERSATIONAL_GREETING":
        return jsonify({"error": "Please ask a specific research question. Greetings are better suited for general chat mode."}), 400

    user = db.session.get(User, user_id)
    now = datetime.now(timezone.utc)
    if user.last_deep_research_timestamp and user.last_deep_research_timestamp.tzinfo is None:
        user.last_deep_research_timestamp = user.last_deep_research_timestamp.replace(tzinfo=timezone.utc)

    if user.last_deep_research_timestamp and (now - user.last_deep_research_timestamp) > timedelta(hours=24):
        user.deep_research_count = 0
        user.last_deep_research_timestamp = None
        db.session.commit()

    if user.deep_research_count >= DEEP_RESEARCH_LIMIT:
        reset_time = user.last_deep_research_timestamp + timedelta(hours=24)
        return jsonify({"error": f"You have reached your daily limit of {DEEP_RESEARCH_LIMIT} reports.", "limit_reached": True, "reset_time": reset_time.isoformat()}), 429

    report_json = {"error": "Failed to generate a valid research report after multiple attempts."}

    for attempt in range(DEEP_RESEARCH_MAX_RETRIES + 1):
        try:
            report_json = generate_deep_research_report(query)
            if "error" not in report_json:
                break
            else:
                print(f"[DEBUG] Attempt {attempt + 1} failed for Deep Research: {report_json.get('error')}. Retrying...")
                time.sleep(DEEP_RESEARCH_RETRY_DELAY)
        except Exception as e:
            print(f"[DEBUG] Exception on attempt {attempt + 1} for Deep Research: {e}. Retrying...")
            time.sleep(DEEP_RESEARCH_RETRY_DELAY)

    if "error" in report_json:
        status_code = report_json.get('status_code', 500)
        return jsonify(report_json), status_code

    try:
        canvas_data = transform_report_to_canvas_data(report_json)

        conversation = get_or_create_conversation(user_id, conversation_id, f"Research: {query[:120]}")
        full_report_message = ChatMessage(
            conversation_id=conversation.id,
            question=f"[DEEP_RESEARCH] {query}",
            answer=json.dumps(report_json),
            sources=report_json.get('references', []),
            canvas_data=canvas_data
        )
        db.session.add(full_report_message)

        abstract_text = report_json.get('abstract', f"The research report on '{query}' was generated successfully.")
        context_message = ChatMessage(
            conversation_id=conversation.id,
            question=f"[REPORT CONTEXT] Abstract for: '{query}'",
            answer=abstract_text,
            sources=None
        )
        db.session.add(context_message)
        conversation.timestamp = datetime.now(timezone.utc)
        if user.deep_research_count == 0:
            user.last_deep_research_timestamp = now
        user.deep_research_count += 1
        db.session.commit()
        return jsonify({
            "success": True,
            "report": report_json,
            "canvas_data": canvas_data,
            "conversation_id": conversation.id,
            "generations_left": DEEP_RESEARCH_LIMIT - user.deep_research_count
        })
    except Exception as e:
        db.session.rollback(); print(f"Error in /deep-research route after report generation: {e}")
        return jsonify({"error": "An internal server error occurred while saving the research results."}), 500

@app.route('/chat', methods=['POST'])
def chat():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    user_query = data['query']
    conversation_id = data.get('conversation_id')
    try:
        history_for_api = get_chat_history_for_api(conversation_id)
        query_type = classify_query(user_query)
        context = ""
        sources_for_frontend = []

        # UPDATED: Conditional logic for query_type
        if query_type == "SEARCH":
            search_results = search_the_web(user_query)
            youtube_results = search_youtube(user_query)
            context_parts = []
            if search_results:
                for item in search_results:
                    source_url = item.get('link', '')
                    sources_for_frontend.append({'url': source_url, 'title': item.get('title', '')})
                    context_parts.append(f"Source [{len(sources_for_frontend)}]: Type: Web Page\nTitle: {item.get('title', '')}\nContent: {item.get('snippet', '')}\nURL: {source_url}")
            if youtube_results:
                for item in youtube_results:
                    video_id = item.get('id', {}).get('videoId')
                    if video_id:
                        video_url = f"https://www.youtube.com/watch?v={video_id}"
                        sources_for_frontend.append({'url': video_url, 'title': item.get('snippet', {}).get('title', '')})
                        context_parts.append(f"Source [{len(sources_for_frontend)}]: Type: YouTube Video\nTitle: {item.get('snippet', {}).get('title', '')}\nURL: {video_url}")
            if context_parts: context = "\n---\n".join(context_parts)
            final_answer = generate_answer(context, user_query, history=history_for_api, use_internal_knowledge=False) # Use context
        elif query_type == "GENERATIVE" or query_type == "CONVERSATIONAL" or query_type == "CONVERSATIONAL_GREETING":
            # For GENERATIVE and CONVERSATIONAL, use AI's internal knowledge without web search context
            final_answer = generate_answer(context="", question=user_query, history=history_for_api, use_internal_knowledge=True)
            sources_for_frontend = [] # No external sources for generative answers
        else: # Fallback, though 'GENERATIVE' should catch most non-search
            final_answer = generate_answer(context="", question=user_query, history=history_for_api, use_internal_knowledge=True)
            sources_for_frontend = []

        conversation = get_or_create_conversation(user_id, conversation_id, user_query[:150])
        new_message = ChatMessage(conversation_id=conversation.id, question=user_query, answer=final_answer, sources=sources_for_frontend)
        db.session.add(new_message)
        conversation.timestamp = datetime.now(timezone.utc)
        db.session.commit()
        return jsonify({"answer": final_answer, "conversation_id": conversation.id, "sources": sources_for_frontend})
    except Exception as e:
        print(f"Error in /chat route: {e}"); db.session.rollback()
        return jsonify({"error": "An internal server error occurred."}), 500

@app.route('/describe-image', methods=['POST'])
def describe_image():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    if 'image' not in request.files: return jsonify({"error": "No image file provided."}), 400
    user_id = session['user']['id']
    image_file = request.files['image']
    prompt_text = request.form.get('prompt', 'Can you describe this image in detail?')
    conversation_id = request.form.get('conversation_id')
    if image_file.filename == '': return jsonify({"error": "No file selected."}), 400
    try:
        description = get_image_description(image_file, prompt_text)
        image_file.seek(0)
        filename = secure_filename(image_file.filename)
        unique_filename = f"{uuid.uuid4()}_{filename}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        image_file.save(file_path)
        web_path = f"/static/uploads/{unique_filename}"
        conversation = get_or_create_conversation(user_id, conversation_id, prompt_text[:150])
        question_content = f"[IMAGE:{web_path}] {prompt_text}"
        new_message = ChatMessage(conversation_id=conversation.id, question=question_content, answer=description)
        db.session.add(new_message)
        conversation.timestamp = datetime.now(timezone.utc)
        db.session.commit()
        return jsonify({"description": description, "conversation_id": conversation.id})
    except Exception as e:
        db.session.rollback()
        print(f"Error in /describe-image route: {e}")
        return jsonify({"error": "An internal server error occurred."}), 500

@app.route('/image-generation-status', methods=['GET'])
def image_generation_status():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    user = db.session.get(User, user_id)
    if not user: return jsonify({"error": "User not found."}), 404
    now = datetime.now(timezone.utc)
    IMAGE_GENERATION_LIMIT = 2
    limit_reached = False
    reset_time_iso = None

    if user.last_generation_timestamp and user.last_generation_timestamp.tzinfo is None:
        user.last_generation_timestamp = user.last_generation_timestamp.replace(tzinfo=timezone.utc)

    if user.last_generation_timestamp and (now - user.last_generation_timestamp) > timedelta(hours=24):
        user.image_generations_count = 0
        user.last_generation_timestamp = None
        db.session.commit()
    generations_left = IMAGE_GENERATION_LIMIT - user.image_generations_count
    if generations_left <= 0:
        limit_reached = True
        if user.last_generation_timestamp:
            reset_time = user.last_generation_timestamp + timedelta(hours=24)
            reset_time_iso = reset_time.isoformat()
    return jsonify({
        "limit_reached": limit_reached,
        "generations_left": generations_left,
        "reset_time": reset_time_iso
    })

@app.route('/generate-image', methods=['POST'])
def generate_image():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    prompt = data.get('prompt')
    conversation_id = data.get('conversation_id')
    if not prompt: return jsonify({"error": "A prompt is required to generate an image."}), 400
    user = db.session.get(User, user_id)
    if not user: return jsonify({"error": "User not found."}), 404
    now = datetime.now(timezone.utc)
    IMAGE_GENERATION_LIMIT = 2
    if user.last_generation_timestamp and user.last_generation_timestamp.tzinfo is None:
        user.last_generation_timestamp = user.last_generation_timestamp.replace(tzinfo=timezone.utc)

    if user.last_generation_timestamp and (now - user.last_generation_timestamp) > timedelta(hours=24):
        user.image_generations_count = 0
        user.last_generation_timestamp = None
        db.session.commit()
    if user.image_generations_count >= IMAGE_GENERATION_LIMIT:
        reset_time = user.last_generation_timestamp + timedelta(hours=24)
        return jsonify({"error": "You have reached your daily image generation limit.", "limit_reached": True, "reset_time": reset_time.isoformat()}), 429

    image_data = generate_image_from_prompt(prompt)
    if isinstance(image_data, str): return jsonify({"error": image_data}), 500
    try:
        unique_filename = f"gen_{uuid.uuid4()}.jpeg"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        with open(file_path, 'wb') as f: f.write(image_data)
        web_path = f"/static/uploads/{unique_filename}"
        conversation = get_or_create_conversation(user_id, conversation_id, prompt[:150])
        answer_content = f"[IMAGE:{web_path}]"
        new_message = ChatMessage(conversation_id=conversation.id, question=prompt, answer=answer_content)
        db.session.add(new_message)
        conversation.timestamp = datetime.now(timezone.utc)
        if user.image_generations_count == 0:
            user.last_generation_timestamp = now
        user.image_generations_count += 1
        db.session.commit()
        return jsonify({
            "success": True, "image_url": web_path,
            "generations_left": IMAGE_GENERATION_LIMIT - user.image_generations_count,
            "conversation_id": conversation.id
        })
    except Exception as e:
        db.session.rollback(); print(f"Error saving generated image: {e}")
        return jsonify({"error": "An internal error occurred while saving the image."}), 500

@app.route('/search-images', methods=['POST'])
def search_images_route():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    query = data.get('query')
    conversation_id = data.get('conversation_id')
    if not query: return jsonify({"error": "A search query is required."}), 400
    try:
        image_results = search_images(query, num_images=10)
        images = [{"link": item.get("link"), "title": item.get("title"), "thumbnail": item.get("image", {}).get("thumbnailLink")} for item in image_results]
        conversation = get_or_create_conversation(user_id, conversation_id, f"Images: {query[:120]}")
        answer_content = json.dumps(images)
        new_message = ChatMessage(conversation_id=conversation.id, question=f"[IMAGE_SEARCH] {query}", answer=answer_content)
        db.session.add(new_message)
        conversation.timestamp = datetime.now(timezone.utc)
        db.session.commit()
        return jsonify({"success": True, "images": images, "conversation_id": conversation.id})
    except Exception as e:
        db.session.rollback(); print(f"Error in /search-images route: {e}")
        return jsonify({"error": "An internal server error occurred."}), 500

# New: API route to add a slide to the presentation
@app.route('/presentation/add_slide', methods=['POST'])
def add_presentation_slide():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    conversation_id = data.get('conversation_id')
    content = data.get('content')
    chat_message_id = data.get('chat_message_id') # Optional: to link to original chat message

    if not conversation_id or not content:
        return jsonify({"error": "Conversation ID and content are required."}), 400

    try:
        # Get the current highest slide order for this conversation
        last_slide = PresentationSlide.query.filter_by(
            user_id=user_id, conversation_id=conversation_id
        ).order_by(PresentationSlide.slide_order.desc()).first()
        new_order = (last_slide.slide_order + 1) if last_slide else 1

        new_slide = PresentationSlide(
            user_id=user_id,
            conversation_id=conversation_id,
            content=content,
            slide_order=new_order,
            chat_message_id=chat_message_id
        )
        db.session.add(new_slide)
        db.session.commit()
        return jsonify({"success": True, "slide": {"id": new_slide.id, "content": new_slide.content, "order": new_slide.slide_order}}), 201
    except Exception as e:
        db.session.rollback()
        print(f"Error adding presentation slide: {e}")
        return jsonify({"error": "An internal server error occurred while adding the slide."}), 500

# New: API route to get all slides for a conversation
@app.route('/presentation/slides/<int:conversation_id>', methods=['GET'])
def get_presentation_slides(conversation_id):
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']

    slides = PresentationSlide.query.filter_by(
        user_id=user_id, conversation_id=conversation_id
    ).order_by(PresentationSlide.slide_order.asc()).all()

    slides_data = [{"id": s.id, "content": s.content, "order": s.slide_order} for s in slides]
    return jsonify(slides_data)

# New: API route to update slide order
@app.route('/presentation/update_slide_order', methods=['POST'])
def update_presentation_slide_order():
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    slides_order = data.get('slides') # Expecting [{id: slide_id, order: new_order}, ...]

    if not slides_order:
        return jsonify({"error": "No slide order data provided."}), 400

    try:
        for slide_data in slides_order:
            slide_id = slide_data.get('id')
            new_order = slide_data.get('order')
            if slide_id is not None and new_order is not None:
                slide = PresentationSlide.query.filter_by(id=slide_id, user_id=user_id).first()
                if slide:
                    slide.slide_order = new_order
        db.session.commit()
        return jsonify({"success": True, "message": "Slide order updated."})
    except Exception as e:
        db.session.rollback()
        print(f"Error updating presentation slide order: {e}")
        return jsonify({"error": "An internal server error occurred while updating slide order."}), 500

# New: API route to delete a slide
@app.route('/presentation/delete_slide/<int:slide_id>', methods=['DELETE'])
def delete_presentation_slide(slide_id):
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']

    try:
        slide = PresentationSlide.query.filter_by(id=slide_id, user_id=user_id).first()
        if not slide:
            return jsonify({"error": "Slide not found or unauthorized access."}), 404

        db.session.delete(slide)
        db.session.commit()

        # Re-order remaining slides to maintain sequential numbering
        remaining_slides = PresentationSlide.query.filter_by(
            user_id=user_id, conversation_id=slide.conversation_id
        ).order_by(PresentationSlide.slide_order.asc()).all()

        for i, s in enumerate(remaining_slides):
            s.slide_order = i + 1
        db.session.commit()

        return jsonify({"success": True, "message": "Slide deleted and re-ordered."})
    except Exception as e:
        db.session.rollback()
        print(f"Error deleting presentation slide: {e}")
        return jsonify({"error": "An internal server error occurred while deleting the slide."}), 500

# New: API route to edit slide content
@app.route('/presentation/edit_slide/<int:slide_id>', methods=['PUT'])
def edit_presentation_slide(slide_id):
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']
    data = request.get_json()
    new_content = data.get('content')

    if not new_content:
        return jsonify({"error": "New content is required."}), 400

    try:
        slide = PresentationSlide.query.filter_by(id=slide_id, user_id=user_id).first()
        if not slide:
            return jsonify({"error": "Slide not found or unauthorized access."}), 404

        slide.content = new_content
        db.session.commit()

        return jsonify({"success": True, "message": "Slide content updated.", "slide": {"id": slide.id, "content": slide.content, "order": slide.slide_order}})
    except Exception as e:
        db.session.rollback()
        print(f"Error editing presentation slide: {e}")
        return jsonify({"error": "An internal server error occurred while editing the slide."}), 500

# New: API route to generate the PPTX presentation
@app.route('/presentation/generate/<int:conversation_id>', methods=['POST'])
def generate_presentation(conversation_id):
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401
    user_id = session['user']['id']

    try:
        slides_data = PresentationSlide.query.filter_by(
            user_id=user_id, conversation_id=conversation_id
        ).order_by(PresentationSlide.slide_order.asc()).all()

        if not slides_data:
            return jsonify({"error": "No slides found for this conversation to generate a presentation."}), 400

        prs = Presentation()
        # Set slide size to 16:9 for a modern look
        prs.slide_width = PptxInches(16)
        prs.slide_height = PptxInches(9)

        # Generate an overall title for the presentation using AI
        overall_title = "AI-Generated Presentation" # Default title
        if slides_data:
            first_slide_content_raw = slides_data[0].content
            model = genai.GenerativeModel(GEMINI_MODEL)
            # Prompt AI to generate a title based on the first slide's content
            title_prompt = f"Given this content which is the first slide's content for a presentation: '{first_slide_content_raw}', generate a captivating and concise overall title for the entire presentation (max 12 words). Ensure the title is professional and summarizes the core theme."
            try:
                response_title = model.generate_content(title_prompt)
                overall_title = response_title.text.strip()
            except Exception as e:
                print(f"Error generating overall presentation title: {e}")
                overall_title = "Synapse AI Presentation"

        # First slide for presentation (Title Slide)
        title_slide_layout = prs.slide_layouts[0] # Title Slide layout is usually index 0
        title_slide = prs.slides.add_slide(title_slide_layout)

        # Title Placeholder
        title_placeholder = title_slide.shapes.title
        if title_placeholder:
            title_placeholder.text = overall_title
            title_text_frame = title_placeholder.text_frame
            p = title_text_frame.paragraphs[0]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            p.alignment = PP_ALIGN.CENTER
        else:
            # Fallback if title placeholder is not found, though rare for layout 0
            left = PptxInches(2)
            top = PptxInches(2)
            width = PptxInches(12)
            height = PptxInches(2)
            title_textbox = title_slide.shapes.add_textbox(left, top, width, height)
            tf = title_textbox.text_frame
            tf.text = overall_title
            p = tf.paragraphs[0]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            p.alignment = PP_ALIGN.CENTER

        # Subtitle Placeholder (e.g., "Presented by Synapse AI Chat")
        try:
            subtitle_placeholder = title_slide.placeholders[1] # Subtitle is usually index 1
            subtitle_placeholder.text = f"Presented by Synapse AI Chat on {datetime.now().strftime('%Y-%m-%d')}"
            subtitle_text_frame = subtitle_placeholder.text_frame
            p = subtitle_text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(128, 128, 128) # Grey
            p.alignment = PP_ALIGN.CENTER
        except IndexError:
            # Fallback if subtitle placeholder is not found
            left = PptxInches(2)
            top = PptxInches(5)
            width = PptxInches(12)
            height = PptxInches(1)
            subtitle_textbox = title_slide.shapes.add_textbox(left, top, width, height)
            tf = subtitle_textbox.text_frame
            tf.text = f"Presented by Synapse AI Chat on {datetime.now().strftime('%Y-%m-%d')}"
            p = tf.paragraphs[0]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(128, 128, 128)
            p.alignment = PP_ALIGN.CENTER


        for i, slide_obj in enumerate(slides_data):
            content = slide_obj.content

            # UPDATED: Generate a more specific image search query from the content
            model = genai.GenerativeModel(GEMINI_MODEL)
            prompt_for_slide_details = f"""
            Given the following content for a presentation slide:
            '{content}'

            Generate a concise and engaging title for this specific slide (max 10 words) and a good, highly relevant image search query to visually represent this slide (max 15 words). Focus the image query on the *core visual concept* of the slide content.
            Respond with a JSON object in this exact format:
            {{
                "slide_title": "Your generated slide title here",
                "image_search_query": "Your generated image search query here"
            }}
            """
            slide_title = f"Slide {i+1}"
            image_search_query = content # Fallback
            try:
                response_details = model.generate_content(prompt_for_slide_details)
                response_text = response_details.text.strip().replace("```json", "").replace("```", "")
                slide_details = json.loads(response_text)
                slide_title = slide_details.get("slide_title", f"Slide {i+1}")
                image_search_query = slide_details.get("image_search_query", content)
            except Exception as e:
                print(f"Error generating slide details with Gemini for content '{content}': {e}")
                # Keep fallback values

            image_found = False
            image_data_io = None

            # Attempt to find and download an image
            image_results = search_images(image_search_query, num_images=1)
            if image_results and image_results[0].get('link'):
                image_url = image_results[0].get('link')
                image_data_io = download_image_for_pptx(image_url)
                if image_data_io:
                    image_found = True

            # Use Title and Content layout (layout 1) for regular slides
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # Assign title to the title placeholder
            title_placeholder = slide.shapes.title
            if title_placeholder:
                title_placeholder.text = slide_title
                title_text_frame = title_placeholder.text_frame
                p = title_text_frame.paragraphs[0]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.theme_color = MSO_THEME_COLOR.DARK_1
                p.alignment = PP_ALIGN.LEFT

            # Find the content placeholder
            body_placeholder = None
            for shape in slide.placeholders:
                # In a 'Title and Content' layout (layout 1), the body placeholder typically has idx=1
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                    body_placeholder = shape
                    break

            if body_placeholder:
                # Resize the content placeholder to make room for an image
                body_placeholder.left = PptxInches(0.5)
                body_placeholder.top = PptxInches(1.5)
                body_placeholder.width = PptxInches(8)
                body_placeholder.height = PptxInches(5.5)

                text_frame = body_placeholder.text_frame
                text_frame.clear() # Clear existing "Click to add text"
                text_frame.word_wrap = True

                # Add content with bullet points and basic formatting
                lines = content.split('\n')
                is_first_line_for_content = True
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue

                    p = text_frame.paragraphs[0] if is_first_line_for_content and len(text_frame.paragraphs) > 0 else text_frame.add_paragraph()

                    if line.startswith('* ') or line.startswith('- '):
                        p.text = line[2:] # Remove bullet char
                        p.level = 1 # Indent for bullet
                        p.font.size = Pt(20)
                        p.font.color.theme_color = MSO_THEME_COLOR.TEXT_2
                    else:
                        p.text = line
                        p.level = 0 # No indent for main paragraphs
                        p.font.size = Pt(20)
                        p.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                    is_first_line_for_content = False

            # Position the image
            if image_found and image_data_io:
                try:
                    # Position image on the right side
                    image_left = PptxInches(9)
                    image_top = PptxInches(1.5)
                    image_height = PptxInches(5.5) # Maintain aspect ratio based on height
                    # The add_picture method will automatically scale the width to maintain aspect ratio
                    slide.shapes.add_picture(image_data_io, image_left, image_top, height=image_height)
                except Exception as img_add_e:
                    print(f"Error adding image to slide: {img_add_e}. Adding placeholder text instead.")
                    # Fallback textbox for image
                    textbox_left = PptxInches(9)
                    textbox_top = PptxInches(1.5)
                    textbox_width = PptxInches(6)
                    textbox_height = PptxInches(5.5)
                    textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
                    text_frame = textbox.text_frame
                    text_frame.text = "Image Not Available"
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif not image_found: # Removed check for image_search_query, always add placeholder if not found
                # Add placeholder text if no image found
                textbox_left = PptxInches(9)
                textbox_top = PptxInches(1.5)
                textbox_width = PptxInches(6)
                textbox_height = PptxInches(5.5)
                textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
                text_frame = textbox.text_frame
                text_frame.text = "Image Not Available"
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


        # Save the presentation
        presentation_title = f"Presentation_{conversation_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        presentation_path = os.path.join(PRESENTATION_OUTPUT_FOLDER, presentation_title)
        prs.save(presentation_path)

        return jsonify({
            "success": True,
            "message": "Presentation generated successfully.",
            "download_url": url_for('download_presentation', filename=presentation_title)
        })

    except Exception as e:
        db.session.rollback()
        print(f"Error generating presentation: {e}")
        return jsonify({"error": f"An internal server error occurred while generating the presentation: {str(e)}"}), 500

@app.route('/download-presentation/<filename>', methods=['GET'])
def download_presentation(filename):
    if 'user' not in session: return jsonify({"error": "Unauthorized"}), 401

    file_path = os.path.join(PRESENTATION_OUTPUT_FOLDER, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "Presentation file not found."}), 404

    return send_file(file_path, as_attachment=True, download_name=filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


@app.after_request
def add_header(response):
    if 'Cache-Control' not in response.headers:
        response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, post-check=0, pre-check=0, max-age=0'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '-1'
    return response

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(host='0.0.0.0', port=8080, debug=True)